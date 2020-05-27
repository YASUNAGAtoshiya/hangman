#coding: utf-8
'''
画像を一気にスライド一枚ずつに貼り付けるコード
ScreenShot *.pngをFigureフォルダに入れると、
フォルダ内の画像を順番通りに一枚ずつ貼り付けたpptファイルを作成する。
'''
import pptx
import os
import re
from pptx import Presentation
from glob import glob

# Presentationインスタンスの作成
ppt = Presentation()
# 幅
width = ppt.slide_width
# 高さ
height = ppt.slide_height

# レイアウト, 0番は白紙
blank_slide_layout = ppt.slide_layouts[0]
# 文字列から数値だけを取り出す
def str2int(str):
    i = int( re.search('\d+',str).group() )
    return i
# 画像ファイルの読み込み
fnms = glob('/Users/yasunagatoshiya/Desktop/Figure/ScreenShot *.png')
# ソート
fnms.sort(key=str2int)
# ファイル毎にループ
for fnm in fnms:
    # 白紙のスライドの追加
    slide = ppt.slides.add_slide(blank_slide_layout)

    # 画像の挿入
    pic = slide.shapes.add_picture(fnm, 0, 0)

    # 中心に移動
    pic.left = int( ( width  - pic.width  ) / 2 )
    pic.top  = int( ( height - pic.height ) / 2 )

# 名前をつけて保存
ppt.save('/Users/yasunagatoshiya/Desktop/Figure/figure.pptx')
